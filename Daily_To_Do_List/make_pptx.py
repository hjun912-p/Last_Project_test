"""
InSIGHT To-Do List & 실험 PPT 생성기
python-pptx 기반, Gamma 스타일 근사
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 팔레트 (Gamma Howlite / Basic Light 근사) ──────────
C_BG       = RGBColor(0xFF, 0xFF, 0xFF)   # 슬라이드 배경
C_TITLE    = RGBColor(0x0D, 0x0D, 0x0D)   # 제목
C_BODY     = RGBColor(0x2D, 0x2D, 0x2D)   # 본문
C_ACCENT   = RGBColor(0x5B, 0x4F, 0xBE)   # 보라 강조 (Gamma 기본색)
C_BOX_BG   = RGBColor(0xEE, 0xEC, 0xF8)   # 카드 배경 (연보라)
C_BOX_B    = RGBColor(0x5B, 0x4F, 0xBE)   # 카드 테두리
C_YELLOW   = RGBColor(0xFF, 0xF3, 0xCD)   # 주의 박스 배경
C_YBOX_B   = RGBColor(0xE5, 0xA8, 0x00)   # 주의 박스 테두리
C_GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)   # 성공기준 박스
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LINE     = RGBColor(0xCC, 0xCC, 0xCC)

W = Inches(13.33)   # 슬라이드 너비 (16:9 기준)
H = Inches(7.5)     # 슬라이드 높이


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # 완전 빈 레이아웃
    slide  = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_BG
    return slide


def add_rect(slide, x, y, w, h, fill=None, border=None, border_w=Pt(1)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1, x, y, w, h)   # 1 = MSO_SHAPE_TYPE.RECTANGLE
    shape.line.color.rgb = border if border else fill if fill else C_LINE
    if not border:
        shape.line.width = Pt(0)
        shape.line.fill.background()
    else:
        shape.line.width = border_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=20, bold=False, color=None, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color if color else C_BODY
    return txBox


def add_textbox_multiline(slide, lines, x, y, w, h,
                           size=16, bold=False, color=None,
                           align=PP_ALIGN.LEFT, line_spacing=None):
    """lines: list of (text, bold_override, size_override, color_override)"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, b, s, c = line, bold, size, color
        else:
            txt = line[0]
            b   = line[1] if len(line) > 1 else bold
            s   = line[2] if len(line) > 2 else size
            c   = line[3] if len(line) > 3 else color

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = txt
        run.font.size  = Pt(s)
        run.font.bold  = b
        run.font.color.rgb = c if c else (color if color else C_BODY)
    return txBox


def top_bar(slide, title_text, subtitle=None):
    """슬라이드 상단 제목 영역"""
    add_text(slide, title_text,
             Inches(0.55), Inches(0.35), Inches(12), Inches(0.7),
             size=36, bold=True, color=C_TITLE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.55), Inches(1.0), Inches(12), Inches(0.4),
                 size=16, color=C_BODY)
    # 하단 구분선
    add_rect(slide, Inches(0.55), Inches(1.35), Inches(12.2), Pt(2),
             fill=C_ACCENT)


def card(slide, text_lines, x, y, w, h,
         bg=None, border=None, title_line=None):
    """카드 박스"""
    bg = bg or C_BOX_BG
    b  = border or C_BOX_B
    add_rect(slide, x, y, w, h, fill=bg, border=b, border_w=Pt(1.2))
    ty = y + Inches(0.15)
    if title_line:
        add_text(slide, title_line,
                 x + Inches(0.2), ty, w - Inches(0.4), Inches(0.4),
                 size=16, bold=True, color=C_ACCENT)
        ty += Inches(0.4)
    add_textbox_multiline(slide, text_lines,
                          x + Inches(0.2), ty,
                          w - Inches(0.4), h - Inches(0.4),
                          size=14, color=C_BODY)


def footer_badge(slide):
    add_text(slide, "InSIGHT · 2026-04-29",
             Inches(9.5), Inches(7.1), Inches(3.5), Inches(0.3),
             size=10, color=RGBColor(0xAA, 0xAA, 0xAA),
             align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
#  TO-DO LIST  (7 슬라이드)
# ════════════════════════════════════════════════════════════

def make_todo():
    prs = new_prs()

    # ── 슬라이드 1: 표지 ─────────────────────────────────────
    sl = blank_slide(prs)
    # 왼쪽 컬러 세로 바
    add_rect(sl, Inches(0), Inches(0), Inches(0.35), H, fill=C_ACCENT)
    add_text(sl, "InSIGHT",
             Inches(0.6), Inches(1.8), Inches(7), Inches(1.2),
             size=72, bold=True, color=C_TITLE)
    add_text(sl, "2026-04-29  팀 오늘의 To Do List",
             Inches(0.6), Inches(3.1), Inches(8), Inches(0.65),
             size=26, bold=True, color=C_TITLE)
    add_text(sl, "AI 생성 콘텐츠 탐지 웹앱 개발팀  ·  이지수 / 박효준 / 진민경 / 김성일 / 신우철",
             Inches(0.6), Inches(3.85), Inches(10), Inches(0.45),
             size=15, color=C_BODY)
    footer_badge(sl)

    # ── 슬라이드 2: 오늘의 배경 ──────────────────────────────
    sl = blank_slide(prs)
    top_bar(sl, "오늘의 배경")
    add_textbox_multiline(sl, [
        ("InSIGHT는 SNS 게시물의 AI 생성 이미지를 자동 판별하는 웹앱입니다.", False, 16, C_BODY),
        ("", False, 8, C_BODY),
        ("현재 1차 필터(C2PA / EXIF / SynthID-FFT)의 탐지 성능이 목표치에 미달하여,", False, 16, C_BODY),
        ("오늘은 마지막 검증과 대안 탐색을 두 트랙으로 병렬 진행합니다.", False, 16, C_BODY),
    ], Inches(0.55), Inches(1.55), Inches(12.2), Inches(1.4))

    card(sl, [
        ("팀장: 이지수  |  팀원: 박효준 · 진민경 · 김성일 · 신우철", False, 15, C_ACCENT),
    ], Inches(0.55), Inches(3.05), Inches(6.5), Inches(0.7),
       bg=C_BOX_BG, border=C_BOX_B)

    add_text(sl, "오늘 핵심 목표",
             Inches(7.5), Inches(3.0), Inches(5.5), Inches(0.45),
             size=18, bold=True, color=C_TITLE)
    add_text(sl, "1차 필터로 채택 가능한 방법 최종 확정",
             Inches(7.5), Inches(3.5), Inches(5.5), Inches(0.5),
             size=15, color=C_BODY)
    footer_badge(sl)

    # ── 슬라이드 3: 오늘의 2대 트랙 ─────────────────────────
    sl = blank_slide(prs)
    top_bar(sl, "오늘의 2대 트랙")

    # Track A 카드
    add_rect(sl, Inches(0.55), Inches(1.6), Inches(5.9), Inches(4.8),
             fill=C_BOX_BG, border=C_BOX_B, border_w=Pt(1.5))
    add_text(sl, "Track A — 대안 탐색",
             Inches(0.75), Inches(1.75), Inches(5.5), Inches(0.5),
             size=20, bold=True, color=C_ACCENT)
    add_text(sl, "담당: 진민경 / 박효준 / 김성일 / 이지수",
             Inches(0.75), Inches(2.25), Inches(5.5), Inches(0.4),
             size=14, bold=False, color=C_BODY)
    add_textbox_multiline(sl, [
        "불가시성 워터마크 검출기 또는",
        "메타데이터 기반 1차 필터 방법 조사",
        "",
        "단순한 정보(워터마크·메타데이터)만으로",
        "분류 가능한 방법 탐색",
        "(ML 모델 개입 없는 결정론적 방법)",
    ], Inches(0.75), Inches(2.75), Inches(5.4), Inches(2.5), size=15)

    # Track B 카드
    add_rect(sl, Inches(6.85), Inches(1.6), Inches(5.9), Inches(4.8),
             fill=C_BOX_BG, border=C_BOX_B, border_w=Pt(1.5))
    add_text(sl, "Track B — 공식 SynthID 검증",
             Inches(7.05), Inches(1.75), Inches(5.5), Inches(0.5),
             size=20, bold=True, color=C_ACCENT)
    add_text(sl, "담당: 신우철",
             Inches(7.05), Inches(2.25), Inches(5.5), Inches(0.4),
             size=14, color=C_BODY)
    add_textbox_multiline(sl, [
        "Vertex AI API KEY로",
        "구글 공식 SynthID 탐지기 연동",
        "",
        "성능·비용 실험으로",
        "1차 필터 적합성 최종 확인",
    ], Inches(7.05), Inches(2.75), Inches(5.4), Inches(2.5), size=15)
    footer_badge(sl)

    # ── 슬라이드 4: Track A 상세 ─────────────────────────────
    sl = blank_slide(prs)
    top_bar(sl, "Track A — 1차 필터 대안 조사",
            "담당: 진민경 / 박효준 / 김성일 / 이지수")

    add_text(sl, "목적: ML 모델 개입 없이 단순 정보(워터마크·메타데이터)만으로 분류 가능한 방법 확정",
             Inches(0.55), Inches(1.55), Inches(12.2), Inches(0.4),
             size=15, color=C_BODY)

    # 조사 항목 박스 2개 (가로 배열, 넓게)
    items = [
        ("불가시성 워터마크",
         "검출 오픈소스 라이브러리\n"
         "- imwatermark, blind-watermark 등\n"
         "- SynthID 외 제3자 워터마크 탐지\n"
         "- 결정론적 신호 기반, API KEY 불필요"),
        ("이미지 메타데이터",
         "기반 AI 탐지 방법\n"
         "- EXIF / XMP / IPTC 심화 분석\n"
         "- C2PA 콘텐츠 서명 파싱\n"
         "- AI 생성 도구 흔적 키워드 탐지"),
    ]
    for i, (t, d) in enumerate(items):
        x = Inches(0.55 + i * 6.2)
        add_rect(sl, x, Inches(2.1), Inches(5.9), Inches(3.3),
                 fill=C_BOX_BG, border=C_BOX_B, border_w=Pt(1.2))
        add_text(sl, t, x + Inches(0.2), Inches(2.25),
                 Inches(5.5), Inches(0.5),
                 size=16, bold=True, color=C_ACCENT)
        add_text(sl, d, x + Inches(0.2), Inches(2.8),
                 Inches(5.5), Inches(2.4), size=14)

    card(sl, [("결과물: 조사 내용 + 후보 방법 팀 채널 공유", True, 14, C_BODY)],
         Inches(0.55), Inches(5.6), Inches(8.0), Inches(0.65),
         bg=C_YELLOW, border=C_YBOX_B)
    footer_badge(sl)

    # ── 슬라이드 5: Track B 상세 ─────────────────────────────
    sl = blank_slide(prs)
    top_bar(sl, "Track B — Vertex AI SynthID 검증 실험",
            "담당: 신우철")

    add_text(sl, "목적: 구글 공식 SynthID 탐지기의 1차 필터 적합성 검증",
             Inches(0.55), Inches(1.55), Inches(12.2), Inches(0.4),
             size=15, color=C_BODY)

    steps = [
        "① Vertex AI API KEY 설정 및 SynthID Watermark Detector 연동",
        "② 탐지 성능 측정: Precision / Recall / F1 / Accuracy",
        "③ 비용 측정: 이미지 1장당 API 호출 비용 (USD)",
        "④ 추론 속도: 1장 평균 처리 시간 (초)",
        "⑤ 1차 필터 목적 적합성 최종 판단 (비용 대비 성능)",
    ]
    for i, s in enumerate(steps):
        y = Inches(2.1 + i * 0.68)
        add_rect(sl, Inches(0.55), y, Inches(11.0), Inches(0.55),
                 fill=C_BOX_BG if i % 2 == 0 else C_WHITE,
                 border=C_LINE, border_w=Pt(0.8))
        add_text(sl, s, Inches(0.75), y + Inches(0.07),
                 Inches(10.6), Inches(0.45), size=15)

    card(sl, [("결과물: 성능 지표 + 비용 리포트 팀 채널 공유", True, 14, C_BODY)],
         Inches(0.55), Inches(5.65), Inches(8.0), Inches(0.65),
         bg=C_YELLOW, border=C_YBOX_B)
    footer_badge(sl)

    # ── 슬라이드 6: 체크리스트 & 성공 기준 ──────────────────
    sl = blank_slide(prs)
    top_bar(sl, "팀 공통 체크리스트 & 오늘의 성공 기준")

    # 왼쪽: 체크리스트
    checks = [
        ("오전", True),
        ("□  Track A 조사 방향 확정 및 시작", False),
        ("□  Track B Vertex AI 연동 및 실험 시작", False),
        ("", False),
        ("오후", True),
        ("□  Track A 조사 결과 팀 채널 공유", False),
        ("□  Track B 실험 결과(지표 + 비용) 공유", False),
        ("", False),
        ("EOD (오늘 마무리 전)", True),
        ("□  1차 필터 채택 방법 최종 결정 논의", False),
        ("□  채택 결과 팀 채널 공유", False),
        ("□  내일 구현 계획 간단 정리", False),
    ]
    add_textbox_multiline(sl,
        [(t, b, 15 if b else 14, C_ACCENT if b else C_BODY) for t, b in checks],
        Inches(0.55), Inches(1.6), Inches(5.8), Inches(5.5))

    # 오른쪽: 성공 기준
    add_text(sl, "오늘의 성공 기준",
             Inches(7.0), Inches(1.6), Inches(6.0), Inches(0.45),
             size=18, bold=True, color=C_TITLE)
    criteria = [
        ("→  1차 필터 채택 가능 방법 1개 이상 확정", C_BODY),
        ("→  Vertex AI SynthID 비용·성능 데이터 확보", C_BODY),
        ("→  팀 전체 방향 통일 완료", C_BODY),
    ]
    for i, (t, c) in enumerate(criteria):
        y = Inches(2.15 + i * 0.75)
        add_rect(sl, Inches(7.0), y, Inches(5.8), Inches(0.6),
                 fill=C_BOX_BG, border=C_BOX_B, border_w=Pt(1))
        add_text(sl, t, Inches(7.2), y + Inches(0.1),
                 Inches(5.4), Inches(0.45), size=14, color=c)

    add_rect(sl, Inches(7.0), Inches(4.55), Inches(5.8), Inches(1.5),
             fill=C_GREEN_BG, border=RGBColor(0x2E, 0x7D, 0x32), border_w=Pt(1))
    add_text(sl, "오늘 두 트랙 모두 완료하면\n내일 1차 필터 구현을 예정대로 시작할 수 있습니다.\n막히는 부분은 즉시 팀 채널에 공유하세요!",
             Inches(7.2), Inches(4.65), Inches(5.4), Inches(1.3),
             size=13, color=RGBColor(0x1B, 0x5E, 0x20))
    footer_badge(sl)

    out = "/Users/woochul/github/Last_Project_test/Daily_To_Do_List/2026.04.29_To_Do_List.pptx"
    prs.save(out)
    print(f"저장 완료: {out}")
    return out


# ════════════════════════════════════════════════════════════
#  실험 PPT  (6 슬라이드)
# ════════════════════════════════════════════════════════════

def make_ppt():
    prs = new_prs()

    # ── 슬라이드 1: 표지 ─────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, Inches(0), Inches(0), W, Inches(0.5), fill=C_ACCENT)
    add_rect(sl, Inches(0), Inches(7.0), W, Inches(0.5), fill=C_ACCENT)
    add_text(sl, "InSIGHT",
             Inches(1.0), Inches(0.9), Inches(11), Inches(1.1),
             size=54, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)
    add_text(sl, "1차 필터 검증 실험",
             Inches(1.0), Inches(2.1), Inches(11), Inches(0.8),
             size=34, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text(sl, "Vertex AI SynthID 탐지기  성능 · 비용 분석",
             Inches(1.0), Inches(2.95), Inches(11), Inches(0.55),
             size=20, color=C_BODY, align=PP_ALIGN.CENTER)
    add_text(sl, "2026-04-29  |  신우철  |  InSIGHT 팀",
             Inches(1.0), Inches(5.8), Inches(11), Inches(0.45),
             size=15, color=C_BODY, align=PP_ALIGN.CENTER)

    # ── 슬라이드 2: 실험 배경 ────────────────────────────────
    sl = blank_slide(prs)
    top_bar(sl, "실험 배경 — 왜 지금 검증하는가")

    problems = [
        ("C2PA",         "SNS 이미지는 대부분 메타데이터 없음 → 사실상 탐지 불가"),
        ("EXIF",         "AI 생성 이미지에 카메라 정보 없음 → inconclusive만 반환"),
        ("SynthID(FFT)", "임계값 기반 근사치 → 실제 워터마크 판독 아님"),
    ]
    for i, (tag, desc) in enumerate(problems):
        y = Inches(1.65 + i * 0.85)
        add_rect(sl, Inches(0.55), y, Inches(1.5), Inches(0.65),
                 fill=C_ACCENT, border=None)
        add_text(sl, tag, Inches(0.6), y + Inches(0.1),
                 Inches(1.4), Inches(0.5), size=14, bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, desc, Inches(2.2), y + Inches(0.12),
                 Inches(10.5), Inches(0.5), size=14, color=C_BODY)

    add_rect(sl, Inches(0.55), Inches(4.4), Inches(12.2), Inches(1.1),
             fill=C_YELLOW, border=C_YBOX_B, border_w=Pt(1.5))
    add_text(sl, "핵심 문제: 1차 필터 무력화 → 모든 이미지가 Gemini 2차 필터로 넘어가 API 비용 과다 발생",
             Inches(0.75), Inches(4.5), Inches(11.8), Inches(0.45),
             size=15, bold=True, color=RGBColor(0x7A, 0x4F, 0x00))
    add_text(sl, "오늘 목적: 구글 공식 SynthID가 1차 필터 역할을 할 수 있는지 비용·성능 양면 검증",
             Inches(0.75), Inches(4.95), Inches(11.8), Inches(0.45),
             size=14, color=C_BODY)
    footer_badge(sl)

    # ── 슬라이드 3: SynthID란? ───────────────────────────────
    sl = blank_slide(prs)
    top_bar(sl, "SynthID란 무엇인가")

    add_text(sl, "Google DeepMind가 개발한 AI 생성 콘텐츠용 불가시성 워터마크 기술",
             Inches(0.55), Inches(1.55), Inches(12.2), Inches(0.45),
             size=16, bold=True, color=C_TITLE)

    left_items = [
        ("핵심 원리", [
            "이미지 픽셀에 인간 눈에 보이지 않는 패턴 삽입",
            "이미지 편집·압축 후에도 워터마크 잔존",
            "Gemini, Imagen 생성 이미지에 자동 적용",
        ]),
        ("탐지 방법", [
            "Vertex AI SynthID Watermark Detector API 호출",
            "탐지 신뢰도 0~1 실수값 반환",
        ]),
    ]
    y_cur = Inches(2.1)
    for title, items in left_items:
        add_text(sl, title, Inches(0.55), y_cur,
                 Inches(6.0), Inches(0.4), size=16, bold=True, color=C_ACCENT)
        y_cur += Inches(0.4)
        for item in items:
            add_text(sl, "  •  " + item, Inches(0.55), y_cur,
                     Inches(6.2), Inches(0.4), size=14)
            y_cur += Inches(0.4)
        y_cur += Inches(0.15)

    add_rect(sl, Inches(7.0), Inches(2.1), Inches(5.8), Inches(3.6),
             fill=C_YELLOW, border=C_YBOX_B, border_w=Pt(1.2))
    add_text(sl, "한계 (알고 있어야 할 것)",
             Inches(7.2), Inches(2.2), Inches(5.4), Inches(0.45),
             size=16, bold=True, color=RGBColor(0x7A, 0x4F, 0x00))
    limits = [
        "구글 AI(Gemini/Imagen) 생성 이미지에만 유효",
        "Midjourney, DALL-E, Stable Diffusion 미적용",
        "공개 오픈소스 SDK 없음",
        "→ Vertex AI API만 가능 (비용 발생)",
    ]
    for i, lim in enumerate(limits):
        add_text(sl, "  •  " + lim,
                 Inches(7.2), Inches(2.75 + i * 0.5),
                 Inches(5.4), Inches(0.45), size=14)
    footer_badge(sl)

    # ── 슬라이드 4: 실험 설계 ────────────────────────────────
    sl = blank_slide(prs)
    top_bar(sl, "실험 설계")

    cols = [
        ("테스트 데이터", [
            "AI 생성 이미지",
            "(Gemini/Imagen 포함 다양한 모델)",
            "",
            "실제 인물 이미지",
            "",
            "기존 벤치마크 데이터셋 활용",
        ]),
        ("측정 지표", [
            "탐지 성능",
            "Precision / Recall / F1 / Accuracy",
            "",
            "비용",
            "이미지 1장당 API 비용 (USD)",
            "",
            "속도",
            "1장 평균 추론 시간 (초)",
        ]),
        ("실험 절차", [
            "① Vertex AI 설정 및 API 연동",
            "② 이미지별 탐지 호출",
            "   + 비용 로그 기록",
            "③ 성능 지표 계산",
            "④ 채택 기준 대비 평가",
            "⑤ 결과 팀 채널 공유",
        ]),
    ]
    for i, (title, items) in enumerate(cols):
        x = Inches(0.55 + i * 4.2)
        add_rect(sl, x, Inches(1.6), Inches(3.9), Inches(4.8),
                 fill=C_BOX_BG, border=C_BOX_B, border_w=Pt(1.2))
        add_text(sl, title, x + Inches(0.2), Inches(1.75),
                 Inches(3.5), Inches(0.45), size=16, bold=True, color=C_ACCENT)
        add_textbox_multiline(sl,
            items, x + Inches(0.2), Inches(2.3),
            Inches(3.5), Inches(3.8), size=14)
    footer_badge(sl)

    # ── 슬라이드 5: 채택 기준 ────────────────────────────────
    sl = blank_slide(prs)
    top_bar(sl, "1차 필터 채택 기준")

    add_text(sl, "채택 조건  (모두 만족해야 함)",
             Inches(0.55), Inches(1.6), Inches(12), Inches(0.45),
             size=18, bold=True, color=C_TITLE)

    criteria = [
        ("Recall  ≥  75 %",      "AI 이미지를 놓치지 않는 것이 핵심 (FN 최소화)"),
        ("장당 비용  ≤  $0.01",  "실용적 수준 — 대량 처리 가능성 확보"),
        ("추론 시간  ≤  3 초",   "사용자 경험 기준"),
    ]
    for i, (crit, reason) in enumerate(criteria):
        y = Inches(2.2 + i * 0.9)
        add_rect(sl, Inches(0.55), y, Inches(3.4), Inches(0.7),
                 fill=C_ACCENT, border=None)
        add_text(sl, crit, Inches(0.65), y + Inches(0.1),
                 Inches(3.2), Inches(0.55), size=15, bold=True,
                 color=C_WHITE)
        add_text(sl, reason, Inches(4.1), y + Inches(0.15),
                 Inches(8.8), Inches(0.45), size=14)

    add_text(sl, "결과별 후속 계획",
             Inches(0.55), Inches(5.05), Inches(12), Inches(0.4),
             size=16, bold=True, color=C_TITLE)
    plans = [
        ("채택",       "insight_app.py 1차 필터에 SynthID API 통합",            C_GREEN_BG),
        ("비용 과다",  "팀원 Track A 대안으로 대체",                              C_YELLOW),
        ("성능 미달",  "앙상블 또는 Gemini 2차 필터 의존 유지",                   C_YELLOW),
    ]
    for i, (tag, plan, bg) in enumerate(plans):
        x = Inches(0.55 + i * 4.2)
        add_rect(sl, x, Inches(5.5), Inches(3.9), Inches(0.85),
                 fill=bg, border=C_LINE, border_w=Pt(0.8))
        add_text(sl, f"[{tag}]  {plan}",
                 x + Inches(0.15), Inches(5.6),
                 Inches(3.6), Inches(0.65), size=13)
    footer_badge(sl)

    # ── 슬라이드 6: 산출물 ───────────────────────────────────
    sl = blank_slide(prs)
    top_bar(sl, "오늘 산출물")

    outputs = [
        ("성능 리포트", [
            "Precision / Recall / F1 / Accuracy 수치",
            "혼동 행렬  (TP / FP / FN / TN)",
            "모델별 탐지율 비교",
            "(Gemini 생성 vs 타사 AI vs 실제)",
        ]),
        ("비용 리포트", [
            "이미지 1장당 평균 비용 (USD)",
            "총 실험 비용",
            "Gemini 2차 필터 대비",
            "비용 효율 비교",
        ]),
        ("최종 권고", [
            "Vertex AI SynthID를",
            "1차 필터로 채택 여부 결론",
            "",
            "팀 채널 공유 →",
            "내일 구현 반영",
        ]),
    ]
    for i, (title, items) in enumerate(outputs):
        x = Inches(0.55 + i * 4.2)
        add_rect(sl, x, Inches(1.6), Inches(3.9), Inches(4.4),
                 fill=C_BOX_BG, border=C_BOX_B, border_w=Pt(1.5))
        add_text(sl, title, x + Inches(0.2), Inches(1.75),
                 Inches(3.5), Inches(0.5), size=17, bold=True, color=C_ACCENT)
        add_textbox_multiline(sl, items,
                              x + Inches(0.2), Inches(2.35),
                              Inches(3.5), Inches(3.3), size=14)
    footer_badge(sl)

    out = "/Users/woochul/github/Last_Project_test/members/woochul/SynthID_검증실험_PPT_20260429.pptx"
    prs.save(out)
    print(f"저장 완료: {out}")
    return out


if __name__ == "__main__":
    make_todo()
    make_ppt()
    print("모두 완료!")
