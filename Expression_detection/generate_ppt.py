"""
PPT 자동 생성 스크립트
실행: python generate_ppt.py
의존성: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import copy

# ── 색상 팔레트 ──────────────────────────────────────────────
C_NAVY   = RGBColor(0x1A, 0x25, 0x4E)   # 진남색 (제목 배경)
C_TEAL   = RGBColor(0x1A, 0x8A, 0x8A)   # 청록 (포인트)
C_ORANGE = RGBColor(0xE6, 0x7E, 0x22)   # 주황 (강조)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xF5, 0xF7, 0xFA)   # 연회색 배경
C_DGRAY  = RGBColor(0x44, 0x44, 0x55)   # 본문 텍스트
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)   # 신규 추가 강조

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # 완전 빈 레이아웃


# ── 헬퍼 함수 ────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=C_DGRAY,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, size=14, bold=False, color=C_DGRAY,
             align=PP_ALIGN.LEFT, space_before=Pt(4)):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return p


def slide_header(slide, title, subtitle=None):
    """상단 헤더 바 + 제목"""
    add_rect(slide, 0, 0, 13.33, 1.2, fill=C_NAVY)
    add_text(slide, title, 0.4, 0.15, 12, 0.7,
             size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.82, 12, 0.35,
                 size=13, bold=False, color=RGBColor(0xAA, 0xCC, 0xDD),
                 align=PP_ALIGN.LEFT)
    # 하단 선
    add_rect(slide, 0, 1.2, 13.33, 0.05, fill=C_TEAL)


def footer(slide, page_num, total=12):
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill=C_NAVY)
    add_text(slide, f'안면 감정 6분류 모델  |  팀장 이지수 · 진민경 · 김성일 · 박효준 · 신우철',
             0.4, 7.12, 10, 0.3, size=9, color=RGBColor(0xAA, 0xCC, 0xDD))
    add_text(slide, f'{page_num} / {total}', 12.0, 7.12, 1.2, 0.3,
             size=9, color=C_WHITE, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
# 슬라이드 1 — 타이틀
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_NAVY)
add_rect(sl, 0, 5.6, 13.33, 1.9, fill=C_TEAL)
add_rect(sl, 0.4, 1.5, 0.12, 2.5, fill=C_ORANGE)

add_text(sl, '안면 감정 6분류 모델 구현', 0.7, 1.4, 12, 1.2,
         size=40, bold=True, color=C_WHITE)
add_text(sl, 'EfficientNetB0 기반 전이학습 + 얼굴 영역 크롭 전처리',
         0.7, 2.7, 12, 0.6, size=20, color=RGBColor(0xAA, 0xDD, 0xEE))
add_text(sl, '비정형 데이터 처리 | 최종 프로젝트', 0.7, 3.4, 8, 0.5,
         size=14, color=RGBColor(0x88, 0xAA, 0xBB))

# 팀원 박스
team = [('팀장', '이지수'), ('팀원', '진민경'), ('팀원', '김성일'),
        ('팀원', '박효준'), ('팀원', '신우철')]
for i, (role, name) in enumerate(team):
    bx = 0.6 + i * 2.4
    add_rect(sl, bx, 5.75, 2.1, 1.3, fill=RGBColor(0x0D, 0x6E, 0x7A))
    add_text(sl, role, bx, 5.8, 2.1, 0.35, size=11,
             color=RGBColor(0xAA, 0xEE, 0xFF), align=PP_ALIGN.CENTER)
    add_text(sl, name, bx, 6.18, 2.1, 0.5, size=18, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(sl, '2026', 11.8, 7.1, 1.3, 0.35, size=10,
         color=RGBColor(0x88, 0xAA, 0xBB), align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
# 슬라이드 2 — 프로젝트 개요
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, '프로젝트 개요', '목표 및 배경')

boxes = [
    ('🎯 목표', '얼굴 이미지로부터 6가지 감정을\n자동으로 분류하는 딥러닝 모델 개발'),
    ('📋 배경', '감정 인식 AI는 의료·교육·\nHCI 등 다양한 분야에 활용 가능'),
    ('📦 데이터셋', 'Face.Sent\n한국인 표정 이미지 데이터셋'),
    ('🏷️ 분류 감정', '기쁨 / 당황 / 분노 / 슬픔\n+ 공포 / 놀람  (총 6종)'),
]
for i, (ttl, body) in enumerate(boxes):
    col = i % 2
    row = i // 2
    bx, by = 0.5 + col * 6.4, 1.6 + row * 2.5
    add_rect(sl, bx, by, 6.0, 2.2, fill=C_WHITE, line=C_TEAL, line_w=Pt(1.5))
    add_rect(sl, bx, by, 6.0, 0.55, fill=C_TEAL)
    add_text(sl, ttl, bx + 0.15, by + 0.07, 5.7, 0.45,
             size=14, bold=True, color=C_WHITE)
    add_text(sl, body, bx + 0.2, by + 0.65, 5.6, 1.4,
             size=15, color=C_DGRAY)

footer(sl, 2)


# ════════════════════════════════════════════════════════════
# 슬라이드 3 — 기존 방식과의 차별점
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, '기존 방식과의 차별점', '무엇이 다른가?')

headers = ['항목', '기존 방식', '✅ 본 프로젝트']
rows = [
    ('감정 클래스', '4종 (기쁨/당황/분노/슬픔)', '6종 (+공포/놀람)'),
    ('모델 아키텍처', 'VGG16, ResNet50', 'EfficientNetB0'),
    ('전처리 방식', '전체 이미지 리사이즈', '얼굴 영역 bbox 크롭'),
    ('학습 전략', 'FC 레이어 단순 학습', '2단계 학습 (Extract → Fine-tune)'),
    ('해석 가능성', '없음', 'Grad-CAM 시각화'),
]

col_w  = [3.0, 4.4, 4.7]
col_x  = [0.35, 3.55, 8.15]
row_h  = 0.72
hdr_y  = 1.45
body_y = 2.17

# 헤더 행
for ci, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    bg = C_NAVY if ci == 0 else (C_TEAL if ci == 1 else RGBColor(0x1A, 0x6E, 0x3A))
    add_rect(sl, cx, hdr_y, cw, 0.6, fill=bg)
    add_text(sl, hdr, cx + 0.1, hdr_y + 0.08, cw - 0.2, 0.45,
             size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 데이터 행
for ri, row_data in enumerate(rows):
    y = body_y + ri * row_h
    bg_alt = RGBColor(0xEB, 0xF5, 0xFB) if ri % 2 == 0 else C_WHITE
    for ci, (val, cx, cw) in enumerate(zip(row_data, col_x, col_w)):
        bg = RGBColor(0xE8, 0xEA, 0xF0) if ci == 0 else bg_alt
        txt_color = C_DGRAY
        if ci == 2:
            txt_color = RGBColor(0x1A, 0x6E, 0x3A)
        add_rect(sl, cx, y, cw, row_h - 0.04, fill=bg,
                 line=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(0.5))
        add_text(sl, val, cx + 0.12, y + 0.12, cw - 0.24, row_h - 0.2,
                 size=12, bold=(ci == 2), color=txt_color, align=PP_ALIGN.CENTER)

footer(sl, 3)


# ════════════════════════════════════════════════════════════
# 슬라이드 4 — 감정 클래스 구성
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, '감정 클래스 구성', 'Ekman(1972) 기본 감정 이론 기반 6종 분류')

emotions = [
    ('기쁨', 'Happy',       '행복하고 즐거운 표정',           RGBColor(0xF3, 0x9C, 0x12), False),
    ('당황', 'Embarrassed', '당황하거나 어리둥절한 표정',       RGBColor(0x9B, 0x59, 0xB6), False),
    ('분노', 'Angry',       '화나거나 짜증나는 표정',           RGBColor(0xE7, 0x4C, 0x3C), False),
    ('슬픔', 'Sad',         '슬프거나 우울한 표정',             RGBColor(0x34, 0x98, 0xDB), False),
    ('공포', 'Scared',      '무섭거나 두려운 표정 ★ 신규 추가', RGBColor(0x1A, 0xBC, 0x9C), True),
    ('놀람', 'Surprised',   '놀라거나 충격받은 표정 ★ 신규 추가', RGBColor(0x27, 0xAE, 0x60), True),
]

for i, (kor, eng, desc, color, is_new) in enumerate(emotions):
    col = i % 3
    row = i // 3
    bx, by = 0.35 + col * 4.28, 1.55 + row * 2.6

    add_rect(sl, bx, by, 4.0, 2.35,
             fill=C_WHITE, line=color, line_w=Pt(2.5))
    add_rect(sl, bx, by, 4.0, 0.65, fill=color)
    add_text(sl, f'{kor}  ({eng})', bx + 0.15, by + 0.1, 3.7, 0.5,
             size=15, bold=True, color=C_WHITE)
    add_text(sl, desc, bx + 0.15, by + 0.78, 3.7, 1.3,
             size=12, color=C_DGRAY)
    if is_new:
        add_rect(sl, bx + 2.7, by + 1.85, 1.15, 0.35,
                 fill=C_GREEN)
        add_text(sl, '신규 추가', bx + 2.72, by + 1.87, 1.1, 0.3,
                 size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

footer(sl, 4)


# ════════════════════════════════════════════════════════════
# 슬라이드 5 — 데이터셋
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, '데이터셋 소개', 'Face.Sent — 한국인 표정 이미지 데이터셋')

items = [
    ('이미지 형식',   'JPG 고해상도 얼굴 사진 (다양한 환경 촬영)'),
    ('JSON 어노테이션', '감정 레이블(faceExp_uploader) + 얼굴 bbox 좌표(minX/Y, maxX/Y)'),
    ('세그멘테이션',  '.npz 마스크 파일 (배경 0 / 머리 1 / 몸 2 / 얼굴 3 / 옷 4 / 기타 5)'),
    ('데이터 분할',   'Train / Validation / Test 3분할'),
    ('다양성',       '성별·연령·직업군 다양, 전문 촬영 + 일반인 참여'),
]

for i, (label, val) in enumerate(items):
    y = 1.55 + i * 1.0
    add_rect(sl, 0.4, y, 3.2, 0.78, fill=C_NAVY)
    add_text(sl, label, 0.55, y + 0.17, 3.0, 0.5,
             size=13, bold=True, color=C_WHITE)
    add_rect(sl, 3.65, y, 9.3, 0.78, fill=C_WHITE,
             line=RGBColor(0xCC, 0xDD, 0xEE), line_w=Pt(1))
    add_text(sl, val, 3.8, y + 0.17, 9.1, 0.5, size=13, color=C_DGRAY)

footer(sl, 5)


# ════════════════════════════════════════════════════════════
# 슬라이드 6 — 전처리 파이프라인
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, '전처리 파이프라인', 'JSON bbox 기반 얼굴 영역 크롭')

# 기존 vs 본 프로젝트 비교
add_rect(sl, 0.3, 1.4, 5.9, 5.3, fill=RGBColor(0xFF, 0xEB, 0xEB),
         line=RGBColor(0xE7, 0x4C, 0x3C), line_w=Pt(1.5))
add_rect(sl, 0.3, 1.4, 5.9, 0.55, fill=RGBColor(0xE7, 0x4C, 0x3C))
add_text(sl, '기존 방식', 0.5, 1.45, 5.5, 0.45,
         size=14, bold=True, color=C_WHITE)

old_steps = ['전체 이미지 로드 (JPG)', '→ 224×224 단순 리사이즈',
             '→ 모델 입력', '',
             '❌ 배경, 몸통 정보 포함', '❌ 불필요한 노이즈 학습']
for i, s in enumerate(old_steps):
    add_text(sl, s, 0.55, 2.05 + i * 0.7, 5.5, 0.6,
             size=13, color=RGBColor(0x7F, 0x1F, 0x1F) if '❌' in s else C_DGRAY)

add_rect(sl, 7.1, 1.4, 5.9, 5.3, fill=RGBColor(0xE8, 0xF8, 0xF1),
         line=C_GREEN, line_w=Pt(1.5))
add_rect(sl, 7.1, 1.4, 5.9, 0.55, fill=C_GREEN)
add_text(sl, '✅ 본 프로젝트', 7.3, 1.45, 5.5, 0.45,
         size=14, bold=True, color=C_WHITE)

new_steps = ['전체 이미지 로드 (JPG)', '→ JSON bbox 좌표 파싱',
             '→ 얼굴 영역 크롭 (+10% 패딩)', '→ 224×224 리사이즈',
             '→ EfficientNet 전처리', '',
             '✅ 배경 노이즈 제거', '✅ 얼굴 특징에 집중']
for i, s in enumerate(new_steps):
    add_text(sl, s, 7.3, 2.05 + i * 0.6, 5.5, 0.55,
             size=13, color=C_GREEN if '✅' in s else C_DGRAY,
             bold='✅' in s)

# 화살표 대체 텍스트
add_text(sl, '→', 6.15, 3.7, 0.8, 0.6,
         size=28, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)

footer(sl, 6)


# ════════════════════════════════════════════════════════════
# 슬라이드 7 — 데이터 증강
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, '데이터 증강 (Data Augmentation)', '과적합 방지 및 일반화 성능 향상')

add_text(sl,
         '얼굴 크롭 후에도 조명·각도·표정 변화에 강건한 모델 학습을 위해 Train 데이터에 다양한 증강을 적용합니다.',
         0.4, 1.35, 12.5, 0.6, size=13, color=C_DGRAY)

augs = [
    ('↻', '회전', '±15도 랜덤 회전'),
    ('↔', '수평 이동', '±8% 이동'),
    ('↕', '수직 이동', '±8% 이동'),
    ('⇔', '좌우 반전', 'Horizontal Flip'),
    ('⊕', '확대/축소', 'Zoom ±10%'),
    ('☀', '밝기 조정', '0.8 ~ 1.2배'),
]

for i, (icon, name, detail) in enumerate(augs):
    col = i % 3
    row = i // 3
    bx, by = 0.35 + col * 4.28, 2.1 + row * 2.3
    add_rect(sl, bx, by, 4.0, 2.05, fill=C_WHITE,
             line=C_TEAL, line_w=Pt(1.5))
    add_text(sl, icon, bx + 0.2, by + 0.2, 0.9, 0.9,
             size=32, bold=True, color=C_TEAL)
    add_text(sl, name, bx + 1.1, by + 0.2, 2.7, 0.5,
             size=15, bold=True, color=C_NAVY)
    add_text(sl, detail, bx + 1.1, by + 0.75, 2.8, 0.5,
             size=13, color=C_DGRAY)

footer(sl, 7)


# ════════════════════════════════════════════════════════════
# 슬라이드 8 — 모델 아키텍처
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, '모델 아키텍처 — EfficientNetB0', '경량·고성능 사전 학습 모델')

# 성능 비교표
add_rect(sl, 0.3, 1.45, 5.8, 0.52, fill=C_NAVY)
for ci, hdr in enumerate(['모델', '파라미터', 'Top-1 정확도']):
    add_text(sl, hdr, 0.4 + ci * 1.93, 1.5, 1.8, 0.4,
             size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

perf = [
    ('VGG16',          '138 M', '71.3%', False),
    ('ResNet50',        ' 25 M', '74.9%', False),
    ('EfficientNetB0 ✓', '  5.3 M', '77.1%', True),
]
for ri, (model_n, param, acc, highlight) in enumerate(perf):
    y = 1.97 + ri * 0.62
    bg = RGBColor(0xE8, 0xF8, 0xF1) if highlight else (
         RGBColor(0xF5, 0xF7, 0xFA) if ri % 2 == 0 else C_WHITE)
    add_rect(sl, 0.3, y, 5.8, 0.58, fill=bg,
             line=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(0.5))
    vals = [model_n, param, acc]
    for ci, v in enumerate(vals):
        add_text(sl, v, 0.4 + ci * 1.93, y + 0.1, 1.8, 0.4,
                 size=12, bold=highlight, color=C_GREEN if highlight else C_DGRAY,
                 align=PP_ALIGN.CENTER)

# 모델 구조 흐름
add_text(sl, '모델 구조 (Input → Output)', 6.5, 1.45, 6.5, 0.45,
         size=14, bold=True, color=C_NAVY)

arch_blocks = [
    ('Input\n224×224×3',       C_DGRAY),
    ('EfficientNetB0\n(Pretrained)',   C_TEAL),
    ('GlobalAvgPooling2D',     C_TEAL),
    ('BatchNormalization',     RGBColor(0x52, 0x7A, 0xAA)),
    ('Dense(512) + Dropout(0.4)', RGBColor(0x8E, 0x44, 0xAD)),
    ('Dense(256) + Dropout(0.3)', RGBColor(0x8E, 0x44, 0xAD)),
    ('Dense(6)  Softmax',      C_ORANGE),
]

bx0 = 6.5
for i, (label, color) in enumerate(arch_blocks):
    by = 1.95 + i * 0.72
    add_rect(sl, bx0, by, 6.2, 0.60, fill=color)
    add_text(sl, label, bx0 + 0.15, by + 0.06, 5.9, 0.52,
             size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if i < len(arch_blocks) - 1:
        add_text(sl, '▼', bx0 + 2.8, by + 0.6, 0.6, 0.12,
                 size=8, color=C_DGRAY, align=PP_ALIGN.CENTER)

footer(sl, 8)


# ════════════════════════════════════════════════════════════
# 슬라이드 9 — 2단계 학습 전략
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, '2단계 학습 전략', 'Feature Extraction → Fine-tuning')

for col, (phase, color_h, color_b, details) in enumerate([
    ('Phase 1\nFeature Extraction',
     C_TEAL, RGBColor(0xE8, 0xF8, 0xF5),
     ['• EfficientNetB0 전체 레이어 동결',
      '• FC 레이어(Dense, Dropout)만 학습',
      '• Learning Rate : 1e-3',
      '• 최대 10 Epochs',
      '• EarlyStopping (patience=5)',
      '• ReduceLROnPlateau (factor=0.5)',
      '• ModelCheckpoint (val_accuracy 기준)']),
    ('Phase 2\nFine-tuning',
     C_ORANGE, RGBColor(0xFE, 0xF9, 0xE7),
     ['• 상위 30개 레이어 학습 허용',
      '• 나머지 레이어는 계속 동결',
      '• Learning Rate : 1e-4  (Phase1의 1/10)',
      '• 최대 10 Epochs',
      '• EarlyStopping (patience=7)',
      '• ReduceLROnPlateau (factor=0.3)',
      '• ModelCheckpoint (val_accuracy 기준)']),
]):
    bx = 0.35 + col * 6.5
    add_rect(sl, bx, 1.4, 6.15, 5.3, fill=color_b,
             line=color_h, line_w=Pt(2))
    add_rect(sl, bx, 1.4, 6.15, 0.75, fill=color_h)
    add_text(sl, phase, bx + 0.2, 1.45, 5.8, 0.65,
             size=16, bold=True, color=C_WHITE)
    for i, d in enumerate(details):
        add_text(sl, d, bx + 0.2, 2.28 + i * 0.6, 5.8, 0.55,
                 size=12, color=C_DGRAY)

footer(sl, 9)


# ════════════════════════════════════════════════════════════
# 슬라이드 10 — Grad-CAM
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, 'Grad-CAM 시각화', '모델 판단 근거 — 어디를 보고 결정했는가?')

add_text(sl,
         'Grad-CAM (Gradient-weighted Class Activation Mapping)',
         0.4, 1.4, 12.5, 0.45, size=16, bold=True, color=C_NAVY)
add_text(sl,
         '예측 클래스에 대한 마지막 Conv 레이어(top_conv)의 그래디언트를 활용해\n'
         '모델이 집중한 이미지 영역을 히트맵(빨강=고집중, 파랑=저집중)으로 표현합니다.',
         0.4, 1.9, 12.5, 0.75, size=13, color=C_DGRAY)

pipeline = [
    ('입력 이미지\n(얼굴 크롭)', C_DGRAY),
    ('EfficientNetB0\nForward Pass', C_TEAL),
    ('top_conv\n그래디언트 계산', RGBColor(0x8E, 0x44, 0xAD)),
    ('가중 평균\n(채널별)', C_ORANGE),
    ('ReLU 적용\n+ 리사이즈', RGBColor(0xE7, 0x4C, 0x3C)),
    ('원본 이미지에\n오버레이', C_GREEN),
]

for i, (label, color) in enumerate(pipeline):
    bx = 0.3 + i * 2.12
    add_rect(sl, bx, 2.85, 1.9, 1.4, fill=color)
    add_text(sl, label, bx + 0.08, 2.95, 1.75, 1.2,
             size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if i < len(pipeline) - 1:
        add_text(sl, '→', bx + 1.92, 3.35, 0.18, 0.5,
                 size=14, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

add_text(sl, '기대 결과 및 의의', 0.4, 4.5, 12.5, 0.45,
         size=14, bold=True, color=C_NAVY)
effects = [
    '• 모델이 눈·입·눈썹 등 감정 핵심 부위에 집중함을 시각적으로 확인',
    '• 예측 오류 시 모델이 어느 영역을 잘못 참조했는지 디버깅 가능',
    '• 모델의 신뢰성과 해석 가능성(Explainability) 제고',
]
for i, e in enumerate(effects):
    add_text(sl, e, 0.4, 5.05 + i * 0.55, 12.5, 0.5, size=13, color=C_DGRAY)

footer(sl, 10)


# ════════════════════════════════════════════════════════════
# 슬라이드 11 — 평가 지표
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, '모델 평가', '정량적 · 정성적 평가 체계')

metrics = [
    ('정확도\n(Accuracy)', 'Test 데이터 전체에 대한\n올바른 예측 비율', C_TEAL),
    ('Confusion Matrix', '클래스별 오분류 패턴\n(절대값 + 정규화)', C_NAVY),
    ('Classification Report', 'Precision / Recall / F1\n클래스별 세부 성능', RGBColor(0x8E, 0x44, 0xAD)),
    ('Grad-CAM', '모델 판단 근거 시각화\n정성적 신뢰성 검증', C_ORANGE),
]

for i, (name, desc, color) in enumerate(metrics):
    col = i % 2
    row = i // 2
    bx, by = 0.35 + col * 6.4, 1.55 + row * 2.55
    add_rect(sl, bx, by, 6.05, 2.3, fill=C_WHITE,
             line=color, line_w=Pt(2))
    add_rect(sl, bx, by, 0.18, 2.3, fill=color)
    add_text(sl, name, bx + 0.35, by + 0.3, 5.5, 0.8,
             size=15, bold=True, color=color)
    add_text(sl, desc, bx + 0.35, by + 1.1, 5.5, 0.9,
             size=13, color=C_DGRAY)

footer(sl, 11)


# ════════════════════════════════════════════════════════════
# 슬라이드 12 — 결론
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_NAVY)
add_rect(sl, 0, 1.1, 13.33, 0.06, fill=C_TEAL)
add_rect(sl, 0, 6.05, 13.33, 0.06, fill=C_ORANGE)

add_text(sl, '결론 및 기대 효과', 0.5, 0.2, 12, 0.75,
         size=30, bold=True, color=C_WHITE)

conclusions = [
    ('✅ 경량·고성능',    'EfficientNetB0(5.3M)로 기존 대비 파라미터 대폭 절감, 정확도 향상'),
    ('✅ 정밀 전처리',    'bbox 기반 얼굴 크롭으로 배경 노이즈 제거, 집중 학습 실현'),
    ('✅ 감정 클래스 확장', '공포·놀람 추가로 실제 감정 표현 6종 커버 (Ekman 이론 기반)'),
    ('✅ 해석 가능성',    'Grad-CAM으로 모델 판단 근거 시각화, 신뢰성 확보'),
]
for i, (title, body) in enumerate(conclusions):
    by = 1.35 + i * 1.05
    add_rect(sl, 0.4, by, 12.4, 0.88, fill=RGBColor(0x0D, 0x1F, 0x44))
    add_text(sl, title, 0.55, by + 0.12, 3.2, 0.65,
             size=13, bold=True, color=C_TEAL)
    add_text(sl, body, 3.9, by + 0.12, 8.7, 0.65,
             size=13, color=RGBColor(0xCC, 0xDD, 0xEE))

add_text(sl, '활용 분야', 0.5, 5.6, 12, 0.4,
         size=14, bold=True, color=C_ORANGE)
apps = ['심리 상담 보조', '온라인 교육\n집중도 모니터링',
        '고객 감정\n분석 서비스', '운전자 위험\n상태 감지']
for i, app in enumerate(apps):
    bx = 0.4 + i * 3.1
    add_rect(sl, bx, 6.1, 2.85, 0.75, fill=RGBColor(0x1A, 0x35, 0x6E))
    add_text(sl, app, bx + 0.1, 6.13, 2.65, 0.65,
             size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

footer(sl, 12)


# ── 저장 ─────────────────────────────────────────────────────
OUTPUT = '안면감정분류_프로젝트발표.pptx'
prs.save(OUTPUT)
print(f'✅ PPT 저장 완료: {OUTPUT}')
