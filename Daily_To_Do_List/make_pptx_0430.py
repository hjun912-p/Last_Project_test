"""
2026-04-30 To-Do List PPTX 생성기
- 1부: KDT 표정 감지 서비스 제출
- 2부: InSIGHT 프로젝트 복귀
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_BG       = RGBColor(0xFF, 0xFF, 0xFF)
C_TITLE    = RGBColor(0x0D, 0x0D, 0x0D)
C_BODY     = RGBColor(0x2D, 0x2D, 0x2D)
C_ACCENT   = RGBColor(0x5B, 0x4F, 0xBE)
C_BOX_BG   = RGBColor(0xEE, 0xEC, 0xF8)
C_BOX_B    = RGBColor(0x5B, 0x4F, 0xBE)
C_YELLOW   = RGBColor(0xFF, 0xF3, 0xCD)
C_YBOX_B   = RGBColor(0xE5, 0xA8, 0x00)
C_GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
C_GREEN_B  = RGBColor(0x2E, 0x7D, 0x32)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LINE     = RGBColor(0xCC, 0xCC, 0xCC)
C_GRAY     = RGBColor(0xAA, 0xAA, 0xAA)
C_RED_BG   = RGBColor(0xFF, 0xEB, 0xEE)
C_RED_B    = RGBColor(0xC6, 0x28, 0x28)

W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_BG
    return slide


def add_rect(slide, x, y, w, h, fill=None, border=None, border_w=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=16, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_BODY
    return tb


def add_multiline(slide, lines, x, y, w, h,
                  size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, b, s, c = line, bold, size, color
        else:
            txt = line[0]
            b   = line[1] if len(line) > 1 else bold
            s   = line[2] if len(line) > 2 else size
            c   = line[3] if len(line) > 3 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.size  = Pt(s)
        run.font.bold  = b
        run.font.color.rgb = c or C_BODY


def top_bar(slide, title, subtitle=None):
    add_text(slide, title,
             Inches(0.55), Inches(0.35), Inches(12), Inches(0.7),
             size=36, bold=True, color=C_TITLE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.55), Inches(1.0), Inches(12), Inches(0.4),
                 size=15, color=C_BODY)
    add_rect(slide, Inches(0.55), Inches(1.35), Inches(12.2), Pt(2), fill=C_ACCENT)


def footer(slide, date="2026-04-30"):
    add_text(slide, f"InSIGHT  ·  {date}",
             Inches(9.5), Inches(7.1), Inches(3.5), Inches(0.3),
             size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
#  슬라이드 생성
# ════════════════════════════════════════════════════════════

def make_todo_0430():
    prs = new_prs()

    # ── 슬라이드 1: 표지 ─────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, Inches(0), Inches(0), Inches(0.35), H, fill=C_ACCENT)

    add_text(sl, "InSIGHT",
             Inches(0.6), Inches(1.5), Inches(7), Inches(1.1),
             size=72, bold=True, color=C_TITLE)
    add_text(sl, "2026-04-30  팀 오늘의 To Do List",
             Inches(0.6), Inches(2.75), Inches(9), Inches(0.65),
             size=26, bold=True, color=C_TITLE)

    # 팀장 부재 알림
    add_rect(sl, Inches(0.6), Inches(3.55), Inches(7.5), Inches(0.65),
             fill=C_RED_BG, border=C_RED_B, border_w=Pt(1.2))
    add_text(sl, "⚠  팀장 이지수님 부재 — 팀원 자율 진행",
             Inches(0.8), Inches(3.65), Inches(7.2), Inches(0.45),
             size=15, bold=True, color=C_RED_B)

    add_text(sl, "진행: 박효준  /  진민경  /  김성일  /  신우철",
             Inches(0.6), Inches(4.35), Inches(9), Inches(0.45),
             size=15, color=C_BODY)
    footer(sl)

    # ── 슬라이드 2: 오늘의 일정 개요 ─────────────────────────
    sl = blank_slide(prs)
    top_bar(sl, "오늘의 일정 개요")

    # 1부 박스
    add_rect(sl, Inches(0.55), Inches(1.6), Inches(5.7), Inches(4.6),
             fill=C_BOX_BG, border=C_BOX_B, border_w=Pt(1.5))
    add_text(sl, "1부  —  KDT 과제 제출",
             Inches(0.75), Inches(1.75), Inches(5.3), Inches(0.5),
             size=20, bold=True, color=C_ACCENT)
    add_text(sl, "우선 완료 후 2부 진행",
             Inches(0.75), Inches(2.25), Inches(5.3), Inches(0.35),
             size=13, italic=True, color=C_GRAY)
    add_multiline(sl, [
        "표정 감지 서비스 개발 및 제출",
        "",
        "  •  폴더: Expression_detection/안면인식자료/",
        "  •  노트북 2종 기반으로 빠르게 완성",
        "  •  KDT 국비지원 제출 기한 준수",
    ], Inches(0.75), Inches(2.7), Inches(5.2), Inches(3.0), size=14)

    # 화살표
    add_text(sl, "→", Inches(6.4), Inches(3.6), Inches(0.6), Inches(0.6),
             size=30, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # 2부 박스
    add_rect(sl, Inches(7.1), Inches(1.6), Inches(5.7), Inches(4.6),
             fill=C_GREEN_BG, border=C_GREEN_B, border_w=Pt(1.5))
    add_text(sl, "2부  —  InSIGHT 복귀",
             Inches(7.3), Inches(1.75), Inches(5.3), Inches(0.5),
             size=20, bold=True, color=C_GREEN_B)
    add_text(sl, "1부 완료 즉시 전환",
             Inches(7.3), Inches(2.25), Inches(5.3), Inches(0.35),
             size=13, italic=True, color=C_GRAY)
    add_multiline(sl, [
        "AI 생성 이미지 탐지 웹앱 작업 재개",
        "",
        "  •  신우철: UI/UX 뼈대 구축",
        "  •  팀원: 각자 담당 작업 재개",
        "  •  진행 상황 팀 채널 공유",
    ], Inches(7.3), Inches(2.7), Inches(5.2), Inches(3.0), size=14)
    footer(sl)

    # ── 슬라이드 3: 1부 상세 — 표정 감지 ────────────────────
    sl = blank_slide(prs)
    top_bar(sl, "1부  —  표정 감지 서비스 (KDT 제출)",
            "Expression_detection/안면인식자료/  기준으로 작업")

    # 파일 목록
    add_text(sl, "작업 파일",
             Inches(0.55), Inches(1.55), Inches(12), Inches(0.4),
             size=16, bold=True, color=C_TITLE)

    files = [
        ("AI_FinalProject_emotion_classification_1.ipynb",
         "표정 분류 모델 1 — 기반 구현"),
        ("AI_FinalProject_emotion_classification_2.ipynb",
         "표정 분류 모델 2 — 개선 / 최종 버전"),
        ("AI개발.비정형_project.guide.26.04.12.pdf",
         "제출 가이드 — 요구사항 및 제출 형식 확인"),
    ]
    for i, (fname, desc) in enumerate(files):
        y = Inches(2.05 + i * 0.95)
        add_rect(sl, Inches(0.55), y, Inches(12.2), Inches(0.8),
                 fill=C_BOX_BG if i % 2 == 0 else C_WHITE,
                 border=C_LINE, border_w=Pt(0.8))
        add_text(sl, fname, Inches(0.75), y + Inches(0.08),
                 Inches(7.0), Inches(0.38), size=13, bold=True, color=C_ACCENT)
        add_text(sl, desc, Inches(7.9), y + Inches(0.18),
                 Inches(4.7), Inches(0.38), size=13, color=C_BODY)

    # 진행 방법
    add_rect(sl, Inches(0.55), Inches(5.0), Inches(12.2), Inches(1.3),
             fill=C_YELLOW, border=C_YBOX_B, border_w=Pt(1.2))
    add_text(sl, "진행 방법",
             Inches(0.75), Inches(5.1), Inches(11.8), Inches(0.38),
             size=15, bold=True, color=RGBColor(0x7A, 0x4F, 0x00))
    add_text(sl,
             "① 가이드 PDF 확인 → 제출 요구사항 파악   "
             "② 노트북 실행 및 결과 확인   "
             "③ 제출 형식에 맞춰 정리   "
             "④ KDT 플랫폼 제출 완료",
             Inches(0.75), Inches(5.55), Inches(11.8), Inches(0.65),
             size=13, color=C_BODY)
    footer(sl)

    # ── 슬라이드 4: 2부 상세 — InSIGHT 복귀 ─────────────────
    sl = blank_slide(prs)
    top_bar(sl, "2부  —  InSIGHT 프로젝트 복귀",
            "1부 완료 즉시 전환 / 팀장 이지수님 부재 시 자율 진행")

    tasks = [
        ("신우철",
         "UI/UX 뼈대 구축\n"
         "• Gradio 기반 메인 화면 레이아웃 설계\n"
         "• 입력/출력 컴포넌트 배치\n"
         "• 모델 선택 버튼 등 인터랙션 구현"),
        ("박효준 / 진민경\n김성일",
         "자율적 자료조사\n"
         "• 2차 필터 구현을 위한 방법 탐색\n"
         "• 1차 / 2차 구분 없이 관심 분야 자유 조사\n"
         "• 조사 결과 팀 채널 공유"),
    ]
    for i, (name, task) in enumerate(tasks):
        x = Inches(0.55 + i * 6.4)
        add_rect(sl, x, Inches(1.6), Inches(6.0), Inches(4.7),
                 fill=C_BOX_BG, border=C_BOX_B, border_w=Pt(1.5))
        add_text(sl, name, x + Inches(0.2), Inches(1.75),
                 Inches(5.6), Inches(0.55), size=18, bold=True, color=C_ACCENT)
        add_multiline(sl,
                      [t for t in task.split("\n")],
                      x + Inches(0.2), Inches(2.4),
                      Inches(5.6), Inches(3.6), size=14)

    add_rect(sl, Inches(0.55), Inches(6.45), Inches(12.2), Inches(0.65),
             fill=C_GREEN_BG, border=C_GREEN_B, border_w=Pt(1))
    add_text(sl, "EOD 목표: 각자 작업 진행 상황 + 이슈 팀 채널 공유  |  이지수 팀장님 복귀 시 즉시 브리핑 가능하도록 기록 유지",
             Inches(0.75), Inches(6.55), Inches(11.8), Inches(0.45),
             size=13, color=RGBColor(0x1B, 0x5E, 0x20))
    footer(sl)

    # ── 슬라이드 5: 체크리스트 ───────────────────────────────
    sl = blank_slide(prs)
    top_bar(sl, "오늘의 체크리스트")

    left = [
        ("1부 — 표정 감지 서비스", True),
        ("□  가이드 PDF 요구사항 확인", False),
        ("□  노트북 1 실행 및 결과 확인", False),
        ("□  노트북 2 실행 및 결과 확인", False),
        ("□  제출 형식 정리 완료", False),
        ("□  KDT 플랫폼 제출 완료", False),
        ("", False),
        ("2부 — InSIGHT 복귀", True),
        ("□  작업 환경 세팅 (python3.12 확인)", False),
        ("□  신우철: UI/UX 뼈대 구현 시작", False),
        ("□  팀원: 2차 필터 관련 자율 자료조사", False),
        ("□  EOD 진행 상황 팀 채널 공유", False),
    ]
    add_multiline(sl,
        [(t, b, 15 if b else 14, C_ACCENT if b else C_BODY) for t, b in left],
        Inches(0.55), Inches(1.6), Inches(5.8), Inches(5.6))

    # 오른쪽: 유의사항
    add_text(sl, "오늘 유의사항",
             Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.45),
             size=18, bold=True, color=C_TITLE)
    notes = [
        ("팀장 부재", "판단이 필요한 결정은 기록만 해두고\n이지수님 복귀 후 확인", C_RED_BG, C_RED_B),
        ("1부 우선", "표정 감지 제출이 오늘 최우선\n완료 전 InSIGHT 작업 시작 금지", C_YELLOW, C_YBOX_B),
        ("기록 유지", "작업 내용·이슈·결정사항\n모두 팀 채널에 남기기", C_GREEN_BG, C_GREEN_B),
    ]
    for i, (tag, desc, bg, bdr) in enumerate(notes):
        y = Inches(2.15 + i * 1.45)
        add_rect(sl, Inches(7.0), y, Inches(5.8), Inches(1.2),
                 fill=bg, border=bdr, border_w=Pt(1.2))
        add_text(sl, tag, Inches(7.2), y + Inches(0.1),
                 Inches(5.4), Inches(0.38), size=14, bold=True, color=bdr)
        add_text(sl, desc, Inches(7.2), y + Inches(0.5),
                 Inches(5.4), Inches(0.6), size=13, color=C_BODY)
    footer(sl)

    out = "/Users/woochul/github/Last_Project_test/Daily_To_Do_List/2026.04.30_To_Do_List.pptx"
    prs.save(out)
    print(f"저장 완료: {out}")
    return out


if __name__ == "__main__":
    make_todo_0430()
